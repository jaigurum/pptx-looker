from pptx import Presentation
from pptx.util import Inches, Pt
import streamlit as st
import tempfile
import os
from PIL import Image
import io
import base64

# Function to convert each section of the PDF into an image
def convert_pdf_sections_to_images(pdf_file):
    try:
        # Open the uploaded PDF using PyMuPDF
        import fitz  # Import here to avoid errors if not used elsewhere
        pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
    except Exception as e:
        st.error(f"Failed to open PDF: {e}")
        return []

    section_images = []

    # Buffer to add at the top of each section
    top_buffer = 20  # You can adjust this value as needed

    # Iterate over each page in the PDF
    for page_num in range(pdf_document.page_count):
        try:
            page = pdf_document.load_page(page_num)  # Load each page
        except Exception as e:
            st.error(f"Failed to load page {page_num}: {e}")
            continue

        # Extract text blocks to identify sections
        try:
            blocks = page.get_text("blocks")  # Extracting blocks of elements
        except Exception as e:
            st.error(f"Failed to extract text blocks from page {page_num}: {e}")
            continue

        sections = []
        section_start = None

        # Identify sections by looking for headings that start with "Section "
        for i, block in enumerate(blocks):
            try:
                # Extract block coordinates and text content
                x0, y0, x1, y1, text = block[:5]
                x0, y0, x1, y1 = float(x0), float(y0), float(x1), float(y1)  # Convert coordinates to float
            except (ValueError, IndexError) as e:
                st.warning(f"Skipping block {i} on page {page_num} due to error: {e}")
                continue

            # Check if the text indicates the start of a new section
            if text and text.strip().startswith("Section "):
                # If we identify a new section heading, mark the end of the previous section
                if section_start is not None:
                    section_end_y = y0  # Use the y-coordinate of the new section as the end of the previous section
                    sections.append((section_start, (0, section_start[1], page.rect.width, section_end_y)))

                # Mark the start of a new section with a top buffer
                section_start = (x0, max(y0 - top_buffer, 0), x1, y1)

        # Capture the last section until the end of the page if it exists
        if section_start is not None:
            section_end = (0, page.rect.height, page.rect.width, page.rect.height)  # End of the page
            sections.append((section_start, section_end))

        # Render each identified section as a separate image
        for idx, section in enumerate(sections):
            try:
                # Extract section coordinates
                (x0_start, y0_start, x1_start, y1_start) = section[0]
                (x0_end, y0_end, x1_end, y1_end) = section[1]

                # Define the rectangle that covers the entire section
                section_rect = fitz.Rect(0, y0_start, page.rect.width, y1_end)
                pix = page.get_pixmap(dpi=150, clip=section_rect)  # Clip to the section

                # Convert the pixmap to a Pillow image
                img_data = pix.tobytes("png")  # Get the image data in PNG format
                image = Image.open(io.BytesIO(img_data))

                # Save the image as a temporary file in a compatible format
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as img_tmp:
                    image.save(img_tmp, format="PNG")
                    img_tmp_path = img_tmp.name

                # Store image path and Pillow image object for later use
                section_images.append((img_tmp_path, image))

            except Exception as e:
                st.error(f"Failed to generate image for section {idx} on page {page_num}: {e}")
                continue

    return section_images

# Function to convert images to PPTX slides
def convert_images_to_pptx(section_images, pdf_filename):
    # Create a new PowerPoint presentation
    try:
        presentation = Presentation()
    except Exception as e:
        st.error(f"Failed to create PowerPoint presentation: {e}")
        return None, None

    # Slide dimensions and margins for fitting images and adding text
    slide_width = Inches(10)
    slide_height = Inches(7.5)

    # Add each image to a separate slide
    for idx, (img_tmp_path, image) in enumerate(section_images):
        try:
            # Add a blank slide to the presentation
            slide_layout = presentation.slide_layouts[6]  # Blank layout
            slide = presentation.slides.add_slide(slide_layout)

            # Resize the image to fit within the slide dimensions while maintaining aspect ratio
            image_width, image_height = image.size
            aspect_ratio = image_width / image_height

            # Calculate new dimensions to fit within slide dimensions
            if aspect_ratio > (slide_width / slide_height):
                new_width = slide_width
                new_height = new_width / aspect_ratio
            else:
                new_height = slide_height
                new_width = new_height * aspect_ratio

            # Center the image on the slide
            x_offset = (slide_width - new_width) / 2
            y_offset = (slide_height - new_height) / 2
            slide.shapes.add_picture(img_tmp_path, x_offset, y_offset, width=new_width, height=new_height)

            # Add a placeholder text box below the image for adding comments or details
            text_box = slide.shapes.add_textbox(Inches(1), slide_height - Inches(1), slide_width - Inches(2), Inches(1))
            text_frame = text_box.text_frame
            text_frame.text = "   "
            text_frame.paragraphs[0].font.size = Pt(14)

        except Exception as e:
            st.error(f"Failed to add image {idx} to slide: {e}")
            continue

    # Save the pptx to a temporary file with the same name as the PDF but with .pptx extension
    try:
        pptx_filename = os.path.splitext(pdf_filename)[0] + ".pptx"
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
            pptx_file = tmp_file.name
            presentation.save(pptx_file)
        return pptx_file, pptx_filename
    except Exception as e:
        st.error(f"Failed to save PowerPoint file: {e}")
        return None, None

# Function to generate a link for automatic download
def create_download_link(file_path, filename):
    with open(file_path, "rb") as f:
        data = f.read()
        b64 = base64.b64encode(data).decode()
        href = f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" download="{filename}">Download PowerPoint</a>'
        return href

# Streamlit app UI
def main():
    st.title("PDF to PPTX Converter")

    # Step 1: Upload PDF
    uploaded_pdf = st.file_uploader("Upload your report as a PDF file", type="pdf")

    if uploaded_pdf is not None:
        # Extract the filename from the uploaded file
        pdf_filename = uploaded_pdf.name

        # Convert PDF sections to images and then convert to PPTX
        with st.spinner("Processing PDF and converting to PPTX..."):
            section_images = convert_pdf_sections_to_images(uploaded_pdf)
            if section_images:
                pptx_path, pptx_filename = convert_images_to_pptx(section_images, pdf_filename)
                if pptx_path:
                    # Automatically initiate the download link
                    download_link = create_download_link(pptx_path, pptx_filename)
                    st.markdown(download_link, unsafe_allow_html=True)
                else:
                    st.error("Failed to convert images to PPTX. Please try again.")

if __name__ == "__main__":
    main()
